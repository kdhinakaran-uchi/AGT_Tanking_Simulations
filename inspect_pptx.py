from pptx import Presentation
from pptx.util import Inches, Pt
import pptx.util

path = r"C:\Users\MBAUser\Downloads\NBA_Draft_Mechanism_Design (1).pptx"
prs = Presentation(path)

print(f"Slide dimensions: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
print(f"Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout.name if slide.slide_layout else "?"
    print(f"--- Slide {i+1} [{layout}] ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  [{shape.shape_type}] {shape.name}: {text[:120]}".encode("ascii", "replace").decode())
        if shape.shape_type == 13:  # picture
            print(f"  [IMAGE] {shape.name}")
    print()
