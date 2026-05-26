"""
Build AGT_Tanking_Deck.pptx from scratch.
Audience: graduate CS students in Algorithmic Game Theory.
Run: python create_deck.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
import copy

DEST = Path("figures/AGT_Tanking_Deck.pptx")
FIGS = Path("figures")

# ── Color palette ─────────────────────────────────────────────────────────────
NAVY      = RGBColor(0x1B, 0x3A, 0x6B)   # primary
RED       = RGBColor(0xE6, 0x39, 0x46)   # accent / NBA red
TEAL      = RGBColor(0x00, 0x89, 0x7B)   # section accent
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # near-black body text
MID_GRAY  = RGBColor(0x55, 0x55, 0x55)
LIGHT     = RGBColor(0xF5, 0xF7, 0xFA)   # slide background tint
PANEL_BG  = RGBColor(0xEE, 0xF2, 0xF7)   # right-panel background

# Figure row colors matching our matplotlib palette
COL_NBA    = RGBColor(0xE6, 0x39, 0x46)
COL_BIL    = RGBColor(0x21, 0x96, 0xF3)
COL_COLA   = RGBColor(0x4C, 0xAF, 0x50)
COL_WL     = RGBColor(0xFF, 0x98, 0x00)
COL_321    = RGBColor(0x9C, 0x27, 0xB0)

# ── Slide geometry ─────────────────────────────────────────────────────────────
SW = Inches(13.33)   # slide width
SH = Inches(7.50)    # slide height

TITLE_H   = Inches(1.05)   # title bar height
BODY_TOP  = Inches(1.15)   # body area starts here
BODY_H    = SH - BODY_TOP - Inches(0.25)
MARGIN_L  = Inches(0.45)
MARGIN_R  = Inches(13.33 - 0.45)

prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank


# ══════════════════════════════════════════════════════════════════════════════
# HELPER PRIMITIVES
# ══════════════════════════════════════════════════════════════════════════════

def add_slide():
    return prs.slides.add_slide(BLANK_LAYOUT)


def rect(slide, left, top, width, height, fill_rgb, alpha=None):
    """Add a solid-fill rectangle with no border."""
    shape = slide.shapes.add_shape(1, left, top, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def tb(slide, left, top, width, height,
       text="", size=18, bold=False, italic=False,
       color=DARK, align=PP_ALIGN.LEFT, wrap=True,
       line_spacing=None):
    """Add a textbox with single-run text."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = None
    p   = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = "Calibri"
    return txb


def tb_para(slide, left, top, width, height,
            lines, size=18, bold=False, color=DARK,
            align=PP_ALIGN.LEFT, line_spacing=1.15,
            heading=None, heading_size=None, heading_color=None):
    """Multi-paragraph textbox.  lines = list of (text, bold_override, size_override)
    or plain strings."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    first = True
    if heading:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = heading
        r.font.bold  = True
        r.font.size  = Pt(heading_size or (size + 2))
        r.font.color.rgb = heading_color or NAVY
        r.font.name  = "Calibri"

    for item in lines:
        if isinstance(item, str):
            text_val, bold_val, size_val = item, bold, size
        else:
            text_val = item[0]
            bold_val = item[1] if len(item) > 1 else bold
            size_val = item[2] if len(item) > 2 else size

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text_val
        r.font.bold   = bold_val
        r.font.size   = Pt(size_val)
        r.font.color.rgb = color
        r.font.name   = "Calibri"

    return txb


def bullet_box(slide, items, left, top, width, height,
               size=17, color=DARK, heading=None, heading_size=20,
               bullet_char="▸", sub_bullet="  –",
               spacing_before=6, line_spacing=1.2):
    """
    Bullet list textbox.
    items can be:
      - str           → regular bullet
      - (str, True)   → bold bullet
      - ("  sub", _)  → indented (starts with spaces)
      - ("---", _)    → blank spacer line
    """
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True

    first = True

    if heading:
        p = tf.paragraphs[0]
        first = False
        r = p.add_run()
        r.text = heading
        r.font.bold  = True
        r.font.size  = Pt(heading_size)
        r.font.color.rgb = NAVY
        r.font.name  = "Calibri"
        p.space_after = Pt(4)

    for item in items:
        if isinstance(item, str):
            text_val, bold_val = item, False
        else:
            text_val, bold_val = item[0], item[1] if len(item) > 1 else False

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False

        if text_val == "---":
            # spacer
            r = p.add_run()
            r.text = " "
            r.font.size = Pt(size // 2)
            continue

        indent = text_val.startswith("   ")
        prefix = (sub_bullet + " ") if indent else (bullet_char + " ")
        display = prefix + text_val.lstrip()

        r = p.add_run()
        r.text = display
        r.font.size  = Pt(size - (2 if indent else 0))
        r.font.bold  = bold_val
        r.font.color.rgb = MID_GRAY if indent else color
        r.font.name  = "Calibri"
        p.space_before = Pt(spacing_before if not indent else 2)

    return txb


def title_bar(slide, title, subtitle=None):
    """Draw the standard navy title bar across the top of a slide."""
    # Background bar
    rect(slide, 0, 0, SW, TITLE_H, NAVY)
    # Red accent stripe on left
    rect(slide, 0, 0, Inches(0.08), TITLE_H, RED)

    # Title text
    tb(slide, Inches(0.22), Inches(0.09), Inches(12.7), Inches(0.65),
       text=title, size=26, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    if subtitle:
        tb(slide, Inches(0.22), Inches(0.65), Inches(12.7), Inches(0.35),
           text=subtitle, size=14, bold=False, color=RGBColor(0xB0, 0xBE, 0xD4),
           align=PP_ALIGN.LEFT)

    return TITLE_H


def section_divider(number, title, subtitle=None):
    """Full-navy section break slide."""
    slide = add_slide()
    rect(slide, 0, 0, SW, SH, NAVY)
    # Decorative red stripe
    rect(slide, 0, Inches(3.2), SW, Inches(0.08), RED)
    # Section label
    tb(slide, Inches(1.0), Inches(1.8), Inches(11.0), Inches(0.6),
       text=f"SECTION {number}", size=14, bold=True,
       color=RGBColor(0x90, 0xA4, 0xBD), align=PP_ALIGN.LEFT)
    tb(slide, Inches(1.0), Inches(2.3), Inches(11.0), Inches(1.2),
       text=title, size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        tb(slide, Inches(1.0), Inches(3.6), Inches(11.0), Inches(0.7),
           text=subtitle, size=18, bold=False,
           color=RGBColor(0xB0, 0xBE, 0xD4), align=PP_ALIGN.LEFT)
    return slide


def results_slide(title, fig_path, takeaways,
                  subtitle=None, chart_w_frac=0.52):
    """Left: chart image.  Right: takeaway bullets on a tinted panel."""
    slide = add_slide()
    title_bar(slide, title, subtitle)

    chart_w = SW * chart_w_frac
    panel_w = SW - chart_w

    # Right panel background
    rect(slide, chart_w, TITLE_H, panel_w, SH - TITLE_H, PANEL_BG)

    # Chart image
    pad = Inches(0.12)
    img_top  = TITLE_H + pad
    img_h    = SH - TITLE_H - 2 * pad
    slide.shapes.add_picture(
        str(fig_path),
        pad, img_top, chart_w - 2 * pad, img_h
    )

    # Takeaways
    bullet_box(slide, takeaways,
               left=chart_w + Inches(0.25),
               top=TITLE_H + Inches(0.30),
               width=panel_w - Inches(0.40),
               height=SH - TITLE_H - Inches(0.4),
               size=17, spacing_before=8, line_spacing=1.3)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
rect(slide, 0, 0, SW, SH, NAVY)

# Decorative block bottom-left
rect(slide, 0, SH - Inches(1.6), Inches(0.5), Inches(1.6), RED)
rect(slide, Inches(0.5), SH - Inches(1.6), Inches(2.2), Inches(1.6),
     RGBColor(0x26, 0x4B, 0x8A))

# Main title
tb(slide, Inches(0.85), Inches(1.4), Inches(11.5), Inches(1.7),
   "Aligning Incentives in the NBA Draft",
   size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Red accent line under title
rect(slide, Inches(0.85), Inches(3.1), Inches(6.5), Inches(0.06), RED)

# Subtitle
tb(slide, Inches(0.85), Inches(3.3), Inches(10.5), Inches(0.75),
   "A Multi-Agent Simulation of Mechanism Design",
   size=22, bold=False, color=RGBColor(0xB0, 0xBE, 0xD4))

# Authors
tb(slide, Inches(0.85), Inches(4.3), Inches(8.0), Inches(0.5),
   "Grant Valentine  &  Kartik Dhinakaran",
   size=18, bold=False, color=RGBColor(0xCC, 0xD6, 0xE4))

# Course / Date
tb(slide, Inches(0.85), Inches(4.85), Inches(8.0), Inches(0.5),
   "Topics in Algorithmic Game Theory  |  May 2026",
   size=14, bold=False, color=RGBColor(0x80, 0x96, 0xAF))

# Key numbers bottom-right
num_x = Inches(10.2)
for i, (num, label) in enumerate([
    ("30", "Teams"), ("82", "Games"), ("50", "Seasons"), ("5", "Mechanisms")
]):
    bx = num_x + i * Inches(0.8)
    tb(slide, bx, Inches(5.8), Inches(0.75), Inches(0.5),
       num, size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
    tb(slide, bx, Inches(6.2), Inches(0.75), Inches(0.4),
       label, size=10, bold=False,
       color=RGBColor(0x90, 0xA4, 0xBD), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 1 DIVIDER: Introduction
# ══════════════════════════════════════════════════════════════════════════════

section_divider(1, "What Is Tanking?",
                "The incentive problem and its history")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: What is tanking?
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "What Is Tanking?")

# Left column – definition
tb(slide, MARGIN_L, Inches(1.25), Inches(5.8), Inches(0.45),
   "The Incentive", size=20, bold=True, color=NAVY)

bullet_box(slide, [
    "Tanking = deliberately losing games to improve future draft position",
    "NBA Draft Lottery: worse record → more ping-pong balls → better odds",
    "---",
    "A team faces a binary season: compete for playoffs, or lose for picks",
    "   When playoff odds fall below a threshold, tanking becomes individually rational",
    "---",
    ("Not a moral failing — a structural game-theory problem", True),
],
left=MARGIN_L, top=Inches(1.75), width=Inches(5.9), height=Inches(4.5),
size=17, spacing_before=7)

# Right column – utility framing
right_x = Inches(7.0)
right_w = Inches(5.9)

rect(slide, right_x - Inches(0.15), TITLE_H + Inches(0.05),
     right_w + Inches(0.2), SH - TITLE_H - Inches(0.15), PANEL_BG)

tb(slide, right_x, Inches(1.25), right_w, Inches(0.45),
   "The Expected-Utility Tradeoff", size=20, bold=True, color=NAVY)

tb(slide, right_x, Inches(1.8), right_w, Inches(0.9),
   "𝔼[U] = p_playoff · V + (1 − p_playoff) · 𝔼[D(k)]",
   size=18, bold=True, color=DARK)

bullet_box(slide, [
    "V = value of a playoff berth (revenue, prestige, competitive stature)",
    "D(k) = value of draft pick #k  (k=1 → $100 pts, falls to 3 pts at k=14)",
    "---",
    "Tanking is rational when:",
    "   ΔE[D(k)] from extra losses  >  V · Δp_playoff lost",
    "---",
    "This tradeoff exists for SOME team under ANY weighted lottery  (Banchio & Munro 2020)",
],
left=right_x, top=Inches(2.8), width=right_w, height=Inches(3.8),
size=16, spacing_before=7)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: History of tanking
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "A Brief History of Tanking",
          subtitle="The incentive has been present since the first draft lottery")

# Timeline entries
events = [
    ("1983", RED,
     "Houston Rockets",
     "Deliberately underperform late-season to secure #1 pick (Ralph Sampson). NBA notices."),
    ("1985", COL_BIL,
     "NBA introduces the lottery",
     "Equal odds for all non-playoff teams. Conspiracy: 'frozen envelope' lets Knicks land Patrick Ewing."),
    ("1990", COL_COLA,
     "Weighted lottery introduced",
     "Worst team gets 16.7% odds. Taylor & Trogdon (2002): eliminated teams lose 2.2× more than expected."),
    ("2013–16", COL_WL,
     "Sam Hinkie's 'The Process'",
     "76ers stockpile picks across 3 seasons. Hinkie forced out. League-wide imitation follows."),
    ("2019", COL_321,
     "NBA flattens top-3 odds to 14% each",
     "Intended to reduce tanking. Instead: displaces incentive to picks #4–7. Tanking continues."),
]

dot_x   = Inches(0.55)
line_x  = Inches(0.75)
entry_x = Inches(1.0)

# Vertical timeline line
rect(slide, line_x, Inches(1.25), Inches(0.03), Inches(5.7), MID_GRAY)

for i, (year, col, name, desc) in enumerate(events):
    y = Inches(1.3 + i * 1.12)
    # Dot
    dot = slide.shapes.add_shape(9, dot_x, y + Inches(0.05),
                                  Inches(0.22), Inches(0.22))  # 9 = oval
    dot.fill.solid()
    dot.fill.fore_color.rgb = col
    dot.line.fill.background()

    # Year
    tb(slide, entry_x, y, Inches(0.7), Inches(0.35),
       year, size=13, bold=True, color=col)
    # Name
    tb(slide, entry_x + Inches(0.75), y, Inches(11.5), Inches(0.35),
       name, size=15, bold=True, color=DARK)
    # Description
    tb(slide, entry_x + Inches(0.75), y + Inches(0.33), Inches(11.3), Inches(0.6),
       desc, size=13, bold=False, color=MID_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: What the NBA has done (lottery reform timeline)
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "NBA Lottery Reforms: A Cycle of Displacement",
          subtitle="Each reform reduces tanking in one place and creates it somewhere else")

reforms = [
    ("1985\nEqual odds",     "All 7 non-playoff teams have identical odds.\nEliminates race-to-bottom, creates race to non-playoffs.",
     COL_COLA),
    ("1990\nWeighted odds",  "Worst team: 16.7% at #1 pick.\nReintroduces tanking incentive, now concentrated at absolute bottom.",
     COL_BIL),
    ("1994\nMore weighting", "Worst team: 25% at #1 pick.\nTanking value even higher for the very worst team.",
     COL_WL),
    ("2019\nFlattened top",  "Top-3 teams each get 14% at #1 pick.\nDisplaces incentive to positions 4–7. Tanking persists.",
     COL_321),
    ("2026\n3-2-1 proposal", "Bottom 3 get fewer balls than teams ranked 4–14.\nOur research tests whether this inverts the incentive.",
     RED),
]

box_w = Inches(2.3)
box_h = Inches(4.5)
gap   = Inches(0.18)
start_x = Inches(0.45)

for i, (title_txt, body_txt, col) in enumerate(reforms):
    bx = start_x + i * (box_w + gap)
    # Top color band
    rect(slide, bx, BODY_TOP, box_w, Inches(0.25), col)
    # Card background
    rect(slide, bx, BODY_TOP + Inches(0.25), box_w, box_h - Inches(0.25), LIGHT)
    # Year / header
    tb(slide, bx + Inches(0.12), BODY_TOP + Inches(0.3), box_w - Inches(0.2),
       Inches(0.65), title_txt, size=14, bold=True, color=NAVY,
       align=PP_ALIGN.CENTER)
    # Body
    tb(slide, bx + Inches(0.12), BODY_TOP + Inches(1.0), box_w - Inches(0.2),
       Inches(3.3), body_txt, size=13, bold=False, color=DARK, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: The impossibility — why fixes keep failing
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Why Reforms Keep Failing: The Banchio Impossibility",
          subtitle="Banchio & Munro (2020), Theorem 1")

# Theorem box
rect(slide, MARGIN_L, Inches(1.2), Inches(12.4), Inches(2.3),
     RGBColor(0xE8, 0xEE, 0xF8))
rect(slide, MARGIN_L, Inches(1.2), Inches(0.07), Inches(2.3), NAVY)

tb(slide, Inches(0.65), Inches(1.3), Inches(12.0), Inches(0.4),
   "Theorem 1 (Banchio & Munro 2020)", size=17, bold=True, color=NAVY)
tb(slide, Inches(0.65), Inches(1.75), Inches(12.0), Inches(1.5),
   "Any draft mechanism that (a) allocates picks based solely on end-of-season standings and "
   "(b) does not treat all non-playoff teams identically, will give some team an incentive "
   "to lose in some possible season history.",
   size=16, bold=False, color=DARK, wrap=True)

# Implications
bullet_box(slide, [
    ("Proof by contradiction:", True),
    "   If mechanism favors worse records and isn't flat, ∃ a rank where one more loss improves lottery odds by more than it costs in playoff probability",
    "   No parameter choice fixes this — it is structural, not parametric",
    "---",
    ("Implication:", True),
    "   No tweak to the current odds table can eliminate tanking",
    "   The only escape routes: (1) lock in standings mid-season, or (2) equal odds (destroys redistribution), or (3) break the standings → lottery link entirely",
],
left=MARGIN_L, top=Inches(3.65), width=Inches(12.4), height=Inches(3.0),
size=16, spacing_before=6)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 2 DIVIDER: Model & Proposals
# ══════════════════════════════════════════════════════════════════════════════

section_divider(2, "The Formal Model & Proposals",
                "Five mechanisms + the math behind agent decisions")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Formal model — high level
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Formal Model: League Setup")

# Left panel
left_w = Inches(6.2)
bullet_box(slide, [
    ("Teams & Skills", True, 18),
    "   n = 30 teams, each with true skill αᵢ ~ LogNormal(0, 0.3)",
    "   Skills evolve: mean reversion + draft pick quality boost",
    "---",
    ("Season Structure", True, 18),
    "   G = 82 games per team (balanced round-robin schedule)",
    "   Top P = 16 teams by wins → playoffs",
    "   Bottom L = 14 teams → draft lottery",
    "---",
    ("Win Probability  (Bradley–Terry with effort floor)", True, 18),
    "   P(A beats B) = (αA · effA) / (αA · effA + αB · effB)",
    "   eff = 0.3 + 0.7 · e  →  even at e = 0, team plays at 30% strength",
    "   Floor δ = 0.3: teams cannot perfectly throw a game",
],
left=MARGIN_L, top=BODY_TOP, width=left_w, height=Inches(5.8),
size=16, spacing_before=5)

# Right panel
right_x = Inches(7.0)
right_w = Inches(5.9)
rect(slide, right_x - Inches(0.1), TITLE_H + Inches(0.05),
     right_w + Inches(0.15), SH - TITLE_H - Inches(0.1), PANEL_BG)

bullet_box(slide, [
    ("Agent Decisions", True, 18),
    "   Effort e ∈ {0, 0.25, 0.5, 0.75, 1.0} (discrete grid)",
    "   Decision at 8 checkpoints per season (every 10 games)",
    "   e < 0.5 = tanking decision",
    "---",
    ("Utility Function", True, 18),
    "   U = V  if team makes playoffs",
    "   U = E[D(k)]  otherwise  (expected pick value from lottery)",
    "---",
    ("Payoff Calibration", True, 18),
    "   V = 200  (playoff >> any single lottery pick)",
    "   D(1)=100, D(2)=65, D(3)=45, ... D(14)=3  pts",
    "   Calibration: consistent with u(playoffs) >> u(picks) (Kazachkov 2020)",
],
left=right_x, top=BODY_TOP, width=right_w, height=Inches(5.8),
size=16, spacing_before=5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: When is tanking rational?
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "When Is Tanking Individually Rational?")

# EU formula
rect(slide, MARGIN_L, Inches(1.2), Inches(12.4), Inches(1.7),
     RGBColor(0xE8, 0xEE, 0xF8))
rect(slide, MARGIN_L, Inches(1.2), Inches(0.07), Inches(1.7), RED)

tb(slide, Inches(0.65), Inches(1.28), Inches(12.0), Inches(0.45),
   "Expected utility comparison (at a checkpoint):", size=15, bold=True, color=NAVY)
tb(slide, Inches(0.65), Inches(1.75), Inches(12.0), Inches(0.95),
   "𝔼[U | compete] = p_playoff(e=1) · V + (1−p_playoff(e=1)) · 𝔼[D(k) | rank_compete]\n"
   "𝔼[U | tank]    = p_playoff(e=0) · V + (1−p_playoff(e=0)) · 𝔼[D(k) | rank_tank]",
   size=15, bold=False, color=DARK, wrap=True)

# Conditions
bullet_box(slide, [
    ("Rational agent tanks when:", True, 18),
    "   The playoff probability with e=1 is negligibly small (team is out of contention)",
    "   AND the lottery rank improvement from e=0 meaningfully raises E[D(k)]",
    "   This is NOT a binary threshold — we compute the normal-approximation p_playoff continuously",
    "---",
    ("Key calibration insight:", True, 18),
    "   With V = 200: E[D(k)] across all lottery slots averages ~24 pts",
    "   Even pick #1 (100 pts at 14% odds) has expected value ≈ 14 pts < < 200",
    "   Tanking is only rational when p_playoff is truly negligible (< ~5%)",
    "---",
    ("Concentration of tanking decisions (our simulation):", True, 18),
    "   76% of tanking decisions: team ranked 21st or worse (deep lottery)",
    "   16% of tanking decisions: team ranked 13–20 (near bubble but behind)",
    "   Teams in playoff position essentially never tank",
],
left=MARGIN_L, top=Inches(3.05), width=Inches(12.4), height=Inches(4.0),
size=16, spacing_before=5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Five mechanisms overview table
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Five Mechanisms Under Study")

# Table using shapes + text
headers = ["Mechanism", "Lottery Basis", "How It Addresses Tanking", "Theoretical Guarantee"]
rows_data = [
    ("NBA Lottery (2019)\n[baseline]",
     "End-of-season standings,\nflattened top-3 at 14%",
     "Doesn't — deliberately loses to\nimprove lottery rank is rational",
     "None (Banchio theorem applies)"),
    ("Bilevel\n(Kazachkov 2020)",
     "Standings at game 70\n(~85% through season)",
     "Locks draft position before\ntanking becomes rational",
     "No post-lock-point incentive;\nresidual pre-lock tanking"),
    ("COLA\n(Highley et al. 2026)",
     "Multi-season accumulated\ntickets (equal/season)",
     "Removes within-season link\nbetween record and lottery odds",
     "Near-full IC: within-season effort\ndoes not affect lottery odds"),
    ("Weighted Loss\n(this paper)",
     "Cumulative loss score;\nearly losses worth more",
     "Decaying reward for tanking;\nno hard cutoff to game",
     "Continuous suppression;\nno strategic discontinuity"),
    ("NBA 3-2-1\n(proposed 2026)",
     "Tier A (worst 3): 2 balls;\nTier B (4–14): 3 balls",
     "Inverts incentive: being\nin bottom 3 is now WORSE",
     "Reduces, does not eliminate;\nTier A guaranteed ≤ pick #12"),
]

col_colors = [NAVY, NAVY, NAVY, NAVY]
col_ws = [Inches(2.1), Inches(2.6), Inches(4.2), Inches(3.8)]
row_h = Inches(0.82)
header_h = Inches(0.48)
start_x = Inches(0.35)
start_y = Inches(1.15)
mech_colors = [COL_NBA, COL_BIL, COL_COLA, COL_WL, COL_321]

# Header row
cx = start_x
for j, (hdr, cw) in enumerate(zip(headers, col_ws)):
    rect(slide, cx, start_y, cw, header_h, NAVY)
    tb(slide, cx + Inches(0.08), start_y + Inches(0.06), cw - Inches(0.1), header_h - Inches(0.1),
       hdr, size=13, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    cx += cw

# Data rows
for i, (row, mc) in enumerate(zip(rows_data, mech_colors)):
    ry = start_y + header_h + i * row_h
    cx = start_x
    bg = LIGHT if i % 2 == 0 else WHITE
    for j, (cell, cw) in enumerate(zip(row, col_ws)):
        rect(slide, cx, ry, cw, row_h,
             RGBColor(0xF0, 0xF4, 0xFF) if (j == 0 and i % 2 == 0) else
             (RGBColor(0xFA, 0xFA, 0xFF) if j == 0 else bg))
        # Mechanism color stripe
        if j == 0:
            rect(slide, cx, ry, Inches(0.06), row_h, mc)
        tb(slide, cx + (Inches(0.13) if j == 0 else Inches(0.08)),
           ry + Inches(0.06), cw - Inches(0.18), row_h - Inches(0.1),
           cell, size=12, bold=(j == 0), color=DARK if j > 0 else NAVY, wrap=True)
        cx += cw


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Bilevel ranking (deep dive)
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Bilevel Ranking  (Kazachkov & Vardi 2020)",
          subtitle="Lock in draft positions before tanking becomes individually rational")

col_w = Inches(5.9)
right_x = Inches(7.1)

bullet_box(slide, [
    ("How it works:", True, 19),
    "   Standings are observed at a preset mid-season breakpoint",
    "   Our calibration: game 70 of 82 (≈ 85% through the season)",
    "   Draft lottery order for non-playoff teams = standings at the lock point",
    "   Post-lock games count for playoff seeding only — not draft order",
    "---",
    ("Why it works:", True, 19),
    "   No team can improve its lottery position after game 70",
    "   Rational agents never tank post-lock: no benefit",
    "   Residual tanking (3.0%) = rational pre-lock positioning",
    "---",
    ("Key result (Kazachkov):", True, 19),
    "   57–72% reduction in tanked games vs. NBA lottery",
    "   Can also IMPROVE competitive balance:",
    "   Late-season tanking distorts standings → bilevel corrects this",
],
left=MARGIN_L, top=BODY_TOP, width=col_w, height=Inches(5.6),
size=16, spacing_before=6)

rect(slide, right_x - Inches(0.1), TITLE_H + Inches(0.05),
     col_w + Inches(0.15), SH - TITLE_H - Inches(0.1), PANEL_BG)

bullet_box(slide, [
    ("Tradeoffs & Critiques:", True, 19),
    "   Breakpoint must be chosen in advance — no principled formula",
    "   Pre-breakpoint tanking incentives remain (teams can game the lock point)",
    "   Teams whose position is locked may rest healthy players (no incentive to compete)",
    "   Fan engagement: confusing when draft order 'freezes' mid-season",
    "---",
    ("In our simulation:", True, 19),
    "   Tanking rate: 3.0%  (vs. 6.9% for NBA lottery)",
    "   Tau distance: 0.398  (honest: 0.403)",
    "   Draft efficiency (weakest quartile): 6.85",
    "   Pre-lock tanking observed: rational agents see value in early-season\n   positioning before the lock point",
],
left=right_x, top=BODY_TOP, width=col_w, height=Inches(5.6),
size=16, spacing_before=6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: COLA
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "COLA: Carry-Over Lottery Allocation  (Highley et al. 2026)",
          subtitle="Full incentive compatibility by breaking the standings → lottery link")

col_w = Inches(5.9)
right_x = Inches(7.1)

bullet_box(slide, [
    ("How it works:", True, 19),
    "   Every non-playoff team receives the same fixed number of tickets each season",
    "   Record does not matter for ticket allocation within a season",
    "   Tickets accumulate across seasons; playoff success / draft wins reduce stockpile",
    "   Lottery draw from cumulative ticket pool",
    "---",
    ("Why it works:", True, 19),
    "   Within any season: effort has no effect on lottery odds",
    "   Rational agent has no incentive to tank",
    "   Long-run: chronically weak teams accumulate tickets → better picks over time",
    "---",
    ("Incentive-compatibility proof:", True, 19),
    "   Diet COLA variant: formally proved IC under assumption that",
    "   making playoffs is always preferred to lottery access",
    "   Classic COLA: simulated IC, not formally proved",
],
left=MARGIN_L, top=BODY_TOP, width=col_w, height=Inches(5.6),
size=16, spacing_before=6)

rect(slide, right_x - Inches(0.1), TITLE_H + Inches(0.05),
     col_w + Inches(0.15), SH - TITLE_H - Inches(0.1), PANEL_BG)

bullet_box(slide, [
    ("Tradeoffs & Critiques:", True, 19),
    "   Multi-season ticket carryover requires complex accounting",
    "   Transition from current system: how to initialize stockpiles fairly?",
    "   IC proof depends on all teams always preferring playoffs",
    "   'Strong draft year' extension uses Bayesian Truth Serum — itself manipulable",
    "---",
    ("In our simulation:", True, 19),
    "   Tanking rate: 0.5%  (near-zero; consistent with IC)",
    "   Tau distance: 0.408  (honest: 0.407 — nearly identical)",
    "   Draft efficiency (weakest quartile): 7.21",
    "   Residual tanking: edge cases where accumulated ticket differentials",
    "   create ambiguous ranking near the lottery threshold",
],
left=right_x, top=BODY_TOP, width=col_w, height=Inches(5.6),
size=16, spacing_before=6)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Weighted Loss & 3-2-1
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Weighted Loss & NBA 3-2-1  (Two Anti-Tanking Designs)")

col_w = Inches(5.9)
right_x = Inches(7.1)

# Left: Weighted Loss
tb(slide, MARGIN_L, Inches(1.15), col_w, Inches(0.42),
   "Weighted Loss Mechanism", size=20, bold=True, color=COL_WL)

bullet_box(slide, [
    "Lottery position = total weighted losses, not final standings",
    "Weight function:  w(g) = 0.5^(g / 20)   (half-life = 20 games)",
    "Early-season loss worth ~4× a late-season loss",
    "---",
    ("Key property:", True),
    "   No hard cutoff — tanking value decays smoothly all season",
    "   Harder for teams to strategically 'time' their tanking window",
    "---",
    ("Simulation results:", True),
    "   Tanking rate: 2.3%",
    "   Residual: rational agents still tank early when decay is slow",
    "   Draft efficiency: 6.92  (solid redistributive performance)",
],
left=MARGIN_L, top=Inches(1.65), width=col_w, height=Inches(5.0),
size=16, spacing_before=6)

# Divider
rect(slide, Inches(6.65), BODY_TOP, Inches(0.04), Inches(5.8), LIGHT)

# Right: NBA 3-2-1
tb(slide, right_x, Inches(1.15), col_w, Inches(0.42),
   "NBA 3-2-1 Proposal", size=20, bold=True, color=COL_321)

bullet_box(slide, [
    "Tier A (worst 3 teams): 2 lottery balls each  →  5.9% odds at #1",
    "Tier B (teams 4–14):    3 lottery balls each  →  8.8% odds at #1",
    "Worst teams guaranteed no worse than pick #12",
    "---",
    ("Key property (inverted incentive):", True),
    "   Being in the bottom 3 is now WORSE than being in positions 4–14",
    "   Rational strategy: avoid the absolute bottom, not target it",
    "---",
    ("Simulation results:", True),
    "   Tanking rate: 0.9%  (strong reduction)",
    "   Draft efficiency: 6.57  — BEST of all mechanisms for weakest teams",
    "   Surprising: fewer balls for worst teams → better outcomes on average",
    "   Fewer extreme lottery misses; value more evenly spread",
],
left=right_x, top=Inches(1.65), width=col_w, height=Inches(5.0),
size=16, spacing_before=6)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 3 DIVIDER: Simulation
# ══════════════════════════════════════════════════════════════════════════════

section_divider(3, "Simulation Structure",
                "What we ran, why, and how agents make decisions")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Experimental conditions table
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Experimental Conditions",
          subtitle="7 conditions × 50 seasons each, all agents decide endogenously")

headers_c = ["Cond.", "Mechanism", "Agent Population", "Purpose"]
rows_c = [
    ("A", "Current NBA (2019)", "All honest (e = 1 always)",
     "Baseline: competitive balance under zero strategic play"),
    ("B", "Current NBA (2019)", "1 rational + 29 honest",
     "Does a single strategic team unravel the honest equilibrium?"),
    ("C", "Current NBA (2019)", "All 30 rational",
     "Equilibrium tanking level under the current system"),
    ("D", "Bilevel  (lock: game 70)", "All 30 rational",
     "Does the bilevel guarantee hold against adaptive agents?"),
    ("E", "COLA", "All 30 rational",
     "Does COLA's incentive compatibility survive strategic agents?"),
    ("F", "Weighted Loss (exp, t½=20)", "All 30 rational",
     "Does continuous decay suppress tanking vs. hard cutoff?"),
    ("G", "NBA 3-2-1", "All 30 rational",
     "Does inverting Tier A odds alter the tanking equilibrium?"),
]
plus = [
    ("★", "Mixed-population sweep",
     "n_rational ∈ {0, 1, 5, 15, 30}",
     "How does tanking scale with rational-agent share?"),
    ("★", "Playoff-value sweep (V)",
     "V ∈ {50, 100, 150, 200, 300}",
     "When do playoff incentives dominate lottery incentives?"),
    ("★", "LLM agents (10 seasons)",
     "All 30 = Claude Haiku",
     "Do LLMs spontaneously discover the tanking incentive?"),
]

col_ws_c = [Inches(0.55), Inches(2.5), Inches(3.3), Inches(6.5)]
header_h_c = Inches(0.42)
row_h_c = Inches(0.55)
start_x_c = Inches(0.3)
start_y_c = Inches(1.15)
cond_colors = [
    RGBColor(0x80, 0x80, 0x80),
    RGBColor(0x80, 0x80, 0x80),
    COL_NBA,
    COL_BIL,
    COL_COLA,
    COL_WL,
    COL_321,
]

# Header
cx = start_x_c
for hdr, cw in zip(headers_c, col_ws_c):
    rect(slide, cx, start_y_c, cw, header_h_c, NAVY)
    tb(slide, cx + Inches(0.06), start_y_c + Inches(0.06),
       cw - Inches(0.1), header_h_c - Inches(0.1),
       hdr, size=13, bold=True, color=WHITE)
    cx += cw

for i, (row, mc) in enumerate(zip(rows_c, cond_colors)):
    ry = start_y_c + header_h_c + i * row_h_c
    bg = LIGHT if i % 2 == 0 else WHITE
    cx = start_x_c
    for j, (cell, cw) in enumerate(zip(row, col_ws_c)):
        rect(slide, cx, ry, cw, row_h_c, bg)
        if j == 0:
            rect(slide, cx, ry, cw, row_h_c, mc)
        fs = 14 if j == 0 else 12
        tb(slide, cx + Inches(0.06), ry + Inches(0.07),
           cw - Inches(0.1), row_h_c - Inches(0.1),
           cell, size=fs, bold=(j == 0),
           color=WHITE if j == 0 else DARK, wrap=True,
           align=PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT)
        cx += cw

# Plus rows (sweeps)
star_y = start_y_c + header_h_c + len(rows_c) * row_h_c + Inches(0.08)
rect(slide, start_x_c, star_y, sum(col_ws_c), Inches(0.3),
     RGBColor(0xD0, 0xD8, 0xF0))
tb(slide, start_x_c + Inches(0.1), star_y + Inches(0.03),
   Inches(8.0), Inches(0.25),
   "Additional sweeps (50 seasons each):",
   size=12, bold=True, color=NAVY)

for i, row in enumerate(plus):
    ry = star_y + Inches(0.3) + i * Inches(0.38)
    bg = LIGHT if i % 2 == 0 else WHITE
    cx = start_x_c
    for j, (cell, cw) in enumerate(zip(row, col_ws_c)):
        rect(slide, cx, ry, cw, Inches(0.38), bg)
        tb(slide, cx + Inches(0.06), ry + Inches(0.04),
           cw - Inches(0.1), Inches(0.32),
           cell, size=12, bold=(j == 0), color=TEAL if j == 0 else DARK,
           align=PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT)
        cx += cw


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Agent design
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Agent Design: Rational vs. LLM",
          subtitle="No hard-coded strategies — behavior emerges from incentives")

col_w = Inches(5.9)
right_x = Inches(7.1)

# Left: Rational
rect(slide, MARGIN_L, BODY_TOP, col_w, Inches(0.38),
     RGBColor(0xE3, 0xED, 0xFB))
tb(slide, MARGIN_L + Inches(0.1), BODY_TOP + Inches(0.05), col_w, Inches(0.3),
   "Rational Agent  (analytical EU maximizer)", size=16, bold=True, color=NAVY)

bullet_box(slide, [
    "At each checkpoint: evaluates EU over effort grid {0, 0.25, 0.5, 0.75, 1.0}",
    "Playoff probability: normal approximation (σ² = n·p·(1−p))",
    "Lottery EV: pre-computed Monte Carlo table (5,000 simulations per lottery structure)",
    "Picks effort with max EU; deterministic given game state",
    "---",
    "Advantage: reproducible, fast, analytically grounded",
    "Advantage: can isolate mechanism effects without LLM noise",
    "Limitation: perfect rationality — not how real GMs reason",
    "Limitation: no learning across seasons (stateless Bayesian agent)",
],
left=MARGIN_L, top=BODY_TOP + Inches(0.42), width=col_w, height=Inches(4.8),
size=15, spacing_before=6)

# Right: LLM
rect(slide, right_x, BODY_TOP, col_w, Inches(0.38),
     RGBColor(0xFB, 0xEC, 0xD8))
tb(slide, right_x + Inches(0.1), BODY_TOP + Inches(0.05), col_w, Inches(0.3),
   "LLM Agent  (Claude Haiku-4-5 via Anthropic API)", size=16, bold=True, color=NAVY)

bullet_box(slide, [
    "Receives structured natural-language prompt at each checkpoint:",
    "   Current record, rank, standings (top 5 / bottom 5)",
    "   Games remaining, playoff gap, lottery tickets",
    "   Mechanism description + full D(k) payoff schedule",
    "Returns JSON:  {\"effort\": float, \"reasoning\": string}",
    "---",
    "System prompt cached per agent instance (reduces API cost ~80%)",
    "No hard-coded tanking rule — behavior fully emergent",
    "---",
    "Advantage: behavioral realism — reasons like a front office",
    "Advantage: spontaneous discovery of game-theoretic incentives",
    "Limitation: expensive (10 seasons of 30 teams × 8 decisions)",
],
left=right_x, top=BODY_TOP + Inches(0.42), width=col_w, height=Inches(4.8),
size=15, spacing_before=6)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 4 DIVIDER: Results
# ══════════════════════════════════════════════════════════════════════════════

section_divider(4, "Results & Findings",
                "50 seasons × 5 mechanisms × rational + honest baselines")


# ══════════════════════════════════════════════════════════════════════════════
# RESULTS SLIDES
# ══════════════════════════════════════════════════════════════════════════════

results_slide(
    title="Result 1: Tanking Rate by Mechanism",
    subtitle="All reforms reduce tanking — but by different amounts and mechanisms",
    fig_path=FIGS / "fig1_tanking_by_mechanism.png",
    takeaways=[
        ("Key numbers (rational agents, 50 seasons):", True),
        "---",
        "  NBA 2019 baseline:   6.9%  ←  ceiling",
        "  Bilevel (lock g=70): 3.0%  ←  pre-lock tanking remains",
        "  Weighted loss:       2.3%  ←  early-season incentive",
        "  NBA 3-2-1:           0.9%  ←  inverted incentive",
        "  COLA:                0.5%  ←  near-zero (near-IC)",
        "---",
        ("Honest agents: 0.0% across all mechanisms", True),
        "---",
        ("Takeaways:", True),
        "All four reforms outperform status quo",
        "No single mechanism perfectly eliminates tanking",
        "COLA and 3-2-1 achieve strongest suppression",
        "Bilevel retains pre-breakpoint incentive",
    ],
)

results_slide(
    title="Result 2: Competitive Balance (Kendall Tau Distance)",
    subtitle="τ = 0: standings perfectly reflect skill  |  τ = 1: fully reversed",
    fig_path=FIGS / "fig2_tau_by_mechanism.png",
    takeaways=[
        ("Tau distance  (rational / honest):", True),
        "---",
        "  NBA 2019:    0.397 / 0.399  (≈ same)",
        "  Bilevel:     0.398 / 0.403",
        "  COLA:        0.408 / 0.407  (identical!)",
        "  Wt. Loss:    0.396 / 0.406",
        "  NBA 3-2-1:   0.409 / 0.407",
        "---",
        ("Surprising finding:", True),
        "With V = 200, rational agents barely change tau vs. honest",
        "Only 6.9% tanking → standings still mostly reflect skill",
        "---",
        ("Draft efficiency (avg pick, weakest quartile):", True),
        "  3-2-1:    6.57  ← best!",
        "  Bilevel:  6.85",
        "  Wt. Loss: 6.92",
        "  COLA:     7.21",
        "  NBA 2019: 7.52  ← worst",
    ],
)

results_slide(
    title="Result 3: Mixed-Population Sweep",
    subtitle="How does tanking scale as the fraction of rational agents grows?",
    fig_path=FIGS / "fig3_mixed_sweep.png",
    takeaways=[
        ("Tanking rate by # rational agents:", True),
        "---",
        "   0 rational:  0.0%",
        "   1 rational:  0.2%",
        "   5 rational:  0.8%",
        "  15 rational:  2.8%",
        "  30 rational:  6.9%",
        "---",
        ("Finding: linear scaling, no cascade", True),
        "Each added rational team ≈ +0.2 pp tanking",
        "Honest equilibrium is NOT fragile at first defection",
        "But unstable over the long run: any front office that",
        "optimizes against the incentive structure gains value",
        "---",
        "Real NBA implication: even if most teams play honestly,",
        "a small number of analytical front offices can rationally tank",
        "without triggering league-wide contagion",
    ],
)

results_slide(
    title="Result 4: Playoff Value Sensitivity",
    subtitle="How does calibrating V (playoff value) change rational tanking?",
    fig_path=FIGS / "fig4_playoff_value_sweep.png",
    takeaways=[
        ("Tanking rate by playoff value V:", True),
        "---",
        "  V =  50:  28.1%  ← high tanking when V is low",
        "  V = 100:  13.1%",
        "  V = 150:   8.8%",
        "  V = 200:   6.9%  ← our baseline",
        "  V = 300:   5.3%",
        "---",
        ("Finding: monotone decline, no sharp threshold", True),
        "Old binary EU model: cliff at V ≈ 50",
        "Probabilistic model: smooth continuous decline",
        "---",
        ("Why this matters:", True),
        "Calibration of V matters enormously for predictions",
        "V = 200 (>> pick value) is consistent with the paper's",
        "core assumption that u(playoffs) >> u(picks)",
        "Real NBA teams likely have V >> 200 in revenue terms",
    ],
)

results_slide(
    title="Result 5: LLM Agents  (Claude Haiku-4-5, 10 seasons)",
    subtitle="Does the tanking incentive emerge spontaneously from language model reasoning?",
    fig_path=FIGS / "fig_llm_comparison.png",
    takeaways=[
        ("LLM: 13.2%  vs.  Rational: 6.9%  (≈ 2×)", True),
        "---",
        ("LLM reasoning patterns:", True),
        '  "At 8-22, playoff probability is negligible (<5%),',
        '   lottery upside exceeds any realistic playoff path"',
        "   → effort = 0.0  (tanking)",
        "---",
        '  "playoff probability is high and 200-pt bonus vastly',
        '   exceeds any lottery pick expected value; compete"',
        "   → effort = 1.0  (competing)",
        "---",
        ("Why LLMs tank more than rational agents:", True),
        "Qualitative elimination heuristics are sharper than",
        "continuous probability-weighted EU calculation",
        "LLMs apply a binary 'out of contention' threshold;",
        "rational agents use a gradient",
        "---",
        "Behavioral finding: real GMs (qualitative reasoners)",
        "may tank closer to LLM rate than rational-agent rate",
    ],
    chart_w_frac=0.46,
)

results_slide(
    title="Result 6: Summary — Tanking vs. Competitive Balance",
    subtitle="No mechanism Pareto-dominates; all involve explicit design tradeoffs",
    fig_path=FIGS / "fig7_summary_dual_axis.png",
    takeaways=[
        ("The tradeoff is real:", True),
        "Mechanisms that suppress tanking most (COLA, 3-2-1)",
        "produce slightly higher tau distance — less tanking means",
        "luck dominates more, so standings reflect skill less",
        "---",
        ("But the effect is small:", True),
        "Tau range across all mechanisms: 0.396 – 0.409",
        "Much smaller than the honest–rational gap in old calibration",
        "V = 200 means most teams compete → standings mostly correct",
        "---",
        ("Policy-relevant insight:", True),
        "There is no free lunch — but the cost is modest",
        "All four reforms substantially reduce tanking",
        "with only marginal impact on competitive balance",
    ],
)


# ══════════════════════════════════════════════════════════════════════════════
# SECTION 5 DIVIDER: Conclusion
# ══════════════════════════════════════════════════════════════════════════════

section_divider(5, "Conclusion",
                "Mechanism comparison, policy recommendations, and open questions")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Mechanism comparison table
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Mechanism Comparison",
          subtitle="Across implementation simplicity, transparency, tanking elimination, and draft efficiency")

def stars(n, total=5):
    return "★" * n + "☆" * (total - n)

comp_headers = [
    "Mechanism", "Simplicity\n(implement)", "Transparency\n(teams/fans)",
    "Tanking\nReduction", "Draft Efficiency\n(weakest quartile)", "Theory"
]
comp_rows = [
    (("NBA 2019\n(baseline)", COL_NBA),
     (stars(5), WHITE, DARK),
     (stars(5), WHITE, DARK),
     ("✗  6.9%", RGBColor(0xFF, 0xEB, 0xEE), RED),
     ("7.52  (worst)", WHITE, DARK),
     ("Impossibility applies", WHITE, MID_GRAY)),

    (("Bilevel\n(Kazachkov)", COL_BIL),
     (stars(3), WHITE, DARK),
     (stars(3), WHITE, DARK),
     ("✓  3.0%", RGBColor(0xE8, 0xF4, 0xFE), COL_BIL),
     ("6.85", WHITE, DARK),
     ("Eliminates post-lock tanking", WHITE, MID_GRAY)),

    (("COLA\n(Highley)", COL_COLA),
     (stars(2), WHITE, DARK),
     (stars(3), WHITE, DARK),
     ("✓✓  0.5%", RGBColor(0xE8, 0xF5, 0xE9), COL_COLA),
     ("7.21", WHITE, DARK),
     ("Near-full IC proof (Diet COLA)", WHITE, MID_GRAY)),

    (("Weighted Loss\n(ours)", COL_WL),
     (stars(3), WHITE, DARK),
     (stars(2), WHITE, DARK),
     ("✓  2.3%", RGBColor(0xFF, 0xF3, 0xE0), COL_WL),
     ("6.92", WHITE, DARK),
     ("Continuous decay, no cutoff", WHITE, MID_GRAY)),

    (("NBA 3-2-1\n(proposed)", COL_321),
     (stars(5), WHITE, DARK),
     (stars(4), WHITE, DARK),
     ("✓✓  0.9%", RGBColor(0xF3, 0xE5, 0xF5), COL_321),
     ("6.57  (BEST ★)", RGBColor(0xF3, 0xE5, 0xF5), COL_321),
     ("Inverts Tier A incentive", WHITE, MID_GRAY)),
]

col_ws_t = [Inches(1.8), Inches(1.5), Inches(1.7), Inches(1.6), Inches(2.0), Inches(4.0)]
header_h_t = Inches(0.55)
row_h_t = Inches(0.88)
start_x_t = Inches(0.25)
start_y_t = Inches(1.15)

# Header
cx = start_x_t
for hdr, cw in zip(comp_headers, col_ws_t):
    rect(slide, cx, start_y_t, cw, header_h_t, NAVY)
    tb(slide, cx + Inches(0.07), start_y_t + Inches(0.05),
       cw - Inches(0.1), header_h_t - Inches(0.08),
       hdr, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER, wrap=True)
    cx += cw

for i, row in enumerate(comp_rows):
    ry = start_y_t + header_h_t + i * row_h_t
    cx = start_x_t
    mech_data = row[0]
    mc = mech_data[1]
    rect(slide, cx, ry, Inches(0.06), row_h_t, mc)
    for j, (cell_data, cw) in enumerate(zip(row, col_ws_t)):
        if isinstance(cell_data, tuple):
            if len(cell_data) == 3:
                cell_text, bg_c, fg_c = cell_data
            else:
                # mechanism column: (text, mech_color) - use white bg, NAVY text
                cell_text = cell_data[0]
                bg_c, fg_c = WHITE, NAVY
        else:
            cell_text, bg_c, fg_c = cell_data, WHITE, DARK
        bg_c_use = RGBColor(0xF8, 0xF8, 0xFF) if bg_c == WHITE and i % 2 == 0 else bg_c
        rect(slide, cx + (Inches(0.06) if j == 0 else 0), ry,
             cw - (Inches(0.06) if j == 0 else 0), row_h_t, bg_c_use)
        align = PP_ALIGN.CENTER if j in (1, 2, 3, 4) else PP_ALIGN.LEFT
        tb(slide, cx + Inches(0.09), ry + Inches(0.08),
           cw - Inches(0.12), row_h_t - Inches(0.12),
           cell_text, size=12 if j != 5 else 11,
           bold=(j == 0), color=fg_c, align=align, wrap=True)
        cx += cw

# Bold recommendation note at bottom
rect(slide, start_x_t, start_y_t + header_h_t + len(comp_rows) * row_h_t + Inches(0.1),
     sum(col_ws_t), Inches(0.38), RGBColor(0xF3, 0xE5, 0xF5))
tb(slide,
   start_x_t + Inches(0.1),
   start_y_t + header_h_t + len(comp_rows) * row_h_t + Inches(0.13),
   sum(col_ws_t) - Inches(0.2), Inches(0.3),
   "★  NBA 3-2-1 achieves the best tanking reduction + draft efficiency combination with highest implementation simplicity",
   size=13, bold=True, color=COL_321)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: What should the NBA do?
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "What Should the NBA Do?",
          subtitle="Our simulation-based policy recommendations")

col_w = Inches(4.0)
gap = Inches(0.15)
bx_h = Inches(5.2)
bx_y = BODY_TOP

for i, (label, color, points) in enumerate([
    ("Near-term: NBA 3-2-1", COL_321, [
        "Easiest to implement — just rebalance balls",
        "Highest transparency — rule is intuitive",
        "0.9% tanking  +  6.57 draft efficiency",
        "No multi-year ramp-up required",
        "Already under serious consideration by NBA",
        "---",
        "Limitation: does not fully eliminate tanking;",
        "restructures rather than removes the incentive",
    ]),
    ("Medium-term: Bilevel", COL_BIL, [
        "Lock standings at game 70 (public, pre-announced)",
        "3.0% tanking — meaningful reduction",
        "Strong theoretical basis (Kazachkov 2020)",
        "No infrastructure overhaul required",
        "---",
        "Limitation: pre-breakpoint incentives remain;",
        "fans need education on why standings 'freeze'",
    ]),
    ("Long-term: COLA", COL_COLA, [
        "Theoretically strongest IC guarantee",
        "0.5% tanking — near-zero",
        "Chronically weak teams accumulate tickets",
        "Formally proved (Diet COLA variant)",
        "---",
        "Limitation: multi-year accounting complexity;",
        "requires transition plan for existing teams;",
        "NBA historically reluctant to restructure",
    ]),
]):
    bx = Inches(0.25) + i * (col_w + gap)
    rect(slide, bx, bx_y, col_w, Inches(0.38), color)
    tb(slide, bx + Inches(0.1), bx_y + Inches(0.05), col_w - Inches(0.15), Inches(0.32),
       label, size=15, bold=True, color=WHITE)
    rect(slide, bx, bx_y + Inches(0.38), col_w, bx_h - Inches(0.38),
         RGBColor(0xF8, 0xF8, 0xFF))
    bullet_box(slide, points,
               left=bx + Inches(0.1),
               top=bx_y + Inches(0.52),
               width=col_w - Inches(0.18),
               height=bx_h - Inches(0.6),
               size=14, spacing_before=5)

# Bottom note
tb(slide, Inches(0.25), bx_y + bx_h + Inches(0.1), Inches(12.8), Inches(0.35),
   "Regardless of mechanism: the Banchio impossibility means we manage, not eliminate, tanking incentives.  "
   "The goal is a regime where the cost of tanking is higher than the benefit for most teams, most of the time.",
   size=12, bold=False, color=MID_GRAY, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Key findings & open questions
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
title_bar(slide, "Key Findings & Open Questions")

left_w = Inches(6.1)
right_x = Inches(6.8)
right_w = Inches(6.1)

tb(slide, MARGIN_L, Inches(1.2), left_w, Inches(0.42),
   "What We Found", size=19, bold=True, color=NAVY)

bullet_box(slide, [
    ("1. All four reforms reduce tanking from the 6.9% NBA baseline", True),
    "   COLA (0.5%) and NBA 3-2-1 (0.9%) achieve strongest suppression",
    "---",
    ("2. Tanking is concentrated far from the playoff bubble", True),
    "   76% of decisions: team ranked 21st or worse",
    "   High V = 200 keeps bubble teams competing",
    "---",
    ("3. Honest equilibrium is stable but not robust", True),
    "   1 rational team → 0.2% aggregate tanking (not contagious)",
    "   But scales linearly: each rational agent adds ~0.2 pp",
    "---",
    ("4. LLM agents tank nearly 2× more than rational agents", True),
    "   Qualitative elimination heuristics are sharper than gradient EU",
    "   Behavioral realism matters for mechanism design predictions",
    "---",
    ("5. NBA 3-2-1 has surprisingly strong draft efficiency", True),
    "   Inverted Tier A incentive improves outcomes for weakest teams",
],
left=MARGIN_L, top=Inches(1.7), width=left_w, height=Inches(5.2),
size=15, spacing_before=5)

rect(slide, right_x - Inches(0.1), TITLE_H + Inches(0.05),
     right_w + Inches(0.15), SH - TITLE_H - Inches(0.1), PANEL_BG)

tb(slide, right_x, Inches(1.2), right_w, Inches(0.42),
   "Open Questions", size=19, bold=True, color=NAVY)

bullet_box(slide, [
    ("Draft class strength", True),
    "We hold D(k) fixed across seasons",
    "In a 'generational draft year', D(1) >> 100 pts",
    "Would mechanisms hold under variable draft quality?",
    "---",
    ("Full 30-team LLM simulation", True),
    "We ran 10 seasons, 1 mechanism",
    "Full 50-season sweep across all mechanisms:",
    "  Do LLMs discover mechanism-specific strategies?",
    "---",
    ("Is some tanking actually optimal?", True),
    "Tau paradox: tanking can IMPROVE standings-skill alignment",
    "If weak teams intentionally lose, draft redistribution works better",
    "Is there a socially optimal tanking rate > 0?",
    "---",
    ("Adversarial robustness", True),
    "What if teams learn each other's strategies over time?",
    "Multi-agent RL rather than static EU maximization",
],
left=right_x, top=Inches(1.7), width=right_w, height=Inches(5.2),
size=15, spacing_before=5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE: Thank you
# ══════════════════════════════════════════════════════════════════════════════

slide = add_slide()
rect(slide, 0, 0, SW, SH, NAVY)
rect(slide, 0, Inches(3.55), SW, Inches(0.07), RED)

# Decorative corner block
rect(slide, SW - Inches(1.2), 0, Inches(1.2), Inches(1.2), RED)
rect(slide, SW - Inches(2.5), 0, Inches(1.3), Inches(1.2),
     RGBColor(0x26, 0x4B, 0x8A))

tb(slide, Inches(1.0), Inches(1.6), Inches(11.0), Inches(1.1),
   "Thank You", size=52, bold=True, color=WHITE)

tb(slide, Inches(1.0), Inches(2.85), Inches(10.0), Inches(0.6),
   "Grant Valentine  &  Kartik Dhinakaran",
   size=20, bold=False, color=RGBColor(0xCC, 0xD6, 0xE4))

tb(slide, Inches(1.0), Inches(3.75), Inches(11.0), Inches(0.55),
   "Questions?", size=26, bold=True, color=RED)

# Summary stats bottom
for i, (num, label) in enumerate([
    ("6.9% → 0.5%", "Tanking range\n(NBA → COLA)"),
    ("2×", "LLM tanks more\nthan rational"),
    ("6.57", "Best draft efficiency\n(NBA 3-2-1)"),
    ("50 seasons", "Simulated per\nmechanism"),
]):
    bx = Inches(1.0 + i * 3.0)
    tb(slide, bx, Inches(4.7), Inches(2.8), Inches(0.55),
       num, size=20, bold=True, color=RED, align=PP_ALIGN.CENTER)
    tb(slide, bx, Inches(5.25), Inches(2.8), Inches(0.55),
       label, size=12, bold=False,
       color=RGBColor(0x90, 0xA4, 0xBD), align=PP_ALIGN.CENTER)

tb(slide, Inches(1.0), Inches(6.5), Inches(11.0), Inches(0.4),
   "Code: github.com/[repo]  |  Data: tanking_sim.db  |  Contact: grantsvalentine@gmail.com",
   size=11, bold=False, color=RGBColor(0x60, 0x78, 0x90))


# ══════════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════════

prs.save(str(DEST))
n = len(prs.slides)
print(f"\nSaved {n}-slide deck -> {DEST}")
print("Slides:")
for i, sl in enumerate(prs.slides, 1):
    # Print the title of each slide (first text box)
    for shape in sl.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t:
                print(f"  {i:2d}. {t[:70]}")
                break
