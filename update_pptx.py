"""
Update NBA_Draft_Mechanism_Design (1).pptx:
- Replace placeholder images on slides 10-14 with our actual figures
- Add pick value chart to slide 9
- Fix stale text notes
- Update LLM reasoning quote with real simulation output
Output: figures/NBA_Draft_Updated.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
import copy, os
from pathlib import Path
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np

SRC  = r"C:\Users\MBAUser\Downloads\NBA_Draft_Mechanism_Design (1).pptx"
DEST = r"C:\Users\MBAUser\Documents\Classes\06. Spring 2026\04. Topics in algorithmic game theory\Final project\AGT_Tanking_Simulations\figures\NBA_Draft_Updated.pptx"
FIGS = Path("figures")

prs = Presentation(SRC)
slides = prs.slides

# ── helper: replace first picture on a slide with a new image ────────────────

def replace_image(slide, img_path, padding_in=0.15):
    """Find the first picture shape, keep its bounding box, swap the image."""
    for shape in slide.shapes:
        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
            left   = shape.left   + Inches(padding_in)
            top    = shape.top    + Inches(padding_in)
            width  = shape.width  - Inches(2 * padding_in)
            height = shape.height - Inches(2 * padding_in)
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(str(img_path), left, top, width, height)
            print(f"    replaced image -> {img_path.name}")
            return True
    return False


def add_image(slide, img_path, left_in, top_in, width_in, height_in):
    slide.shapes.add_picture(
        str(img_path),
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in),
    )
    print(f"    added image -> {img_path.name}")


def remove_shape_by_text(slide, text_fragment):
    """Remove a text shape whose text contains text_fragment."""
    for shape in list(slide.shapes):
        if shape.has_text_frame:
            full = " ".join(p.text for p in shape.text_frame.paragraphs)
            if text_fragment in full:
                shape._element.getparent().remove(shape._element)
                print(f"    removed shape containing: '{text_fragment[:60]}'")
                return True
    return False


def update_text(slide, old_fragment, new_text):
    """Replace the text of the first shape containing old_fragment."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full = " ".join(p.text for p in shape.text_frame.paragraphs)
            if old_fragment in full:
                tf = shape.text_frame
                for para in tf.paragraphs:
                    for run in para.runs:
                        if old_fragment in run.text:
                            run.text = run.text.replace(old_fragment, new_text)
                            print(f"    updated text: '{old_fragment[:50]}' -> '{new_text[:50]}'")
                            return True
    return False


# ── make LLM comparison figure ───────────────────────────────────────────────

def make_llm_comparison():
    COLORS = {
        "nba_lottery": "#E63946",
        "llm":         "#795548",
    }
    fig, ax = plt.subplots(figsize=(7, 4.5))
    labels = ["Rational Agents\n(NBA Lottery)", "LLM Agents\n(Claude Haiku)"]
    rates  = [15.0, 44.2]
    colors = [COLORS["nba_lottery"], COLORS["llm"]]

    bars = ax.bar([0, 1], rates, color=colors, width=0.5,
                  edgecolor="white", linewidth=0.5)
    for bar, r in zip(bars, rates):
        ax.text(bar.get_x() + bar.get_width()/2, r + 0.8,
                f"{r:.1f}%", ha="center", va="bottom",
                fontsize=14, fontweight="bold")

    ax.set_xticks([0, 1])
    ax.set_xticklabels(labels, fontsize=12)
    ax.set_ylabel("Tanking rate (%)", fontsize=12)
    ax.set_ylim(0, 55)
    ax.set_title("LLM Agents Tank 3× More Than Rational Agents\n(NBA Lottery, 10 seasons)",
                 fontsize=13, fontweight="bold")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    ax.annotate("Qualitative heuristics\ntrigger sharper\ntanking response",
                xy=(1, 44.2), xytext=(1.35, 35),
                arrowprops=dict(arrowstyle="->", color="gray"),
                fontsize=10, color="gray", ha="left")

    plt.tight_layout()
    out = FIGS / "fig_llm_comparison.png"
    plt.savefig(out, bbox_inches="tight", dpi=150)
    plt.close()
    print(f"  created {out.name}")
    return out


llm_fig = make_llm_comparison()

# ── slide-by-slide updates ────────────────────────────────────────────────────

# Slide 3 (index 2): remove stale "Need to add the latest mechanism (3-2-1)" note
print("Slide 3: removing stale note...")
remove_shape_by_text(slides[2], "Need to add the latest mechanism")

# Slide 4 (index 3): update "Four proposals" -> "Five proposals"
print("Slide 4: updating proposal count...")
update_text(slides[3], "Four proposals claim", "Five proposals claim")

# Slide 9 (index 8): add pick value chart in the diagram area
print("Slide 9: adding pick value chart...")
# The slide has a hand-drawn kink diagram; place our chart on the right half
add_image(slides[8], FIGS / "fig5_pick_value_schedule.png",
          left_in=5.2, top_in=1.1, width_in=4.6, height_in=3.8)

# Slide 10 (index 9): playoff value sensitivity
print("Slide 10: replacing image with fig4...")
if not replace_image(slides[9], FIGS / "fig4_playoff_value_sweep.png"):
    add_image(slides[9], FIGS / "fig4_playoff_value_sweep.png",
              left_in=0.5, top_in=1.2, width_in=6.5, height_in=3.8)
remove_shape_by_text(slides[9], "League tanking rates by playoff value")

# Slide 11 (index 10): tanking rate by mechanism
print("Slide 11: replacing image with fig1...")
if not replace_image(slides[10], FIGS / "fig1_tanking_by_mechanism.png"):
    add_image(slides[10], FIGS / "fig1_tanking_by_mechanism.png",
              left_in=0.3, top_in=1.2, width_in=7.0, height_in=3.8)
remove_shape_by_text(slides[10], "Tanking rates by proposal")

# Slide 12 (index 11): Kendall tau
print("Slide 12: replacing image with fig2...")
if not replace_image(slides[11], FIGS / "fig2_tau_by_mechanism.png"):
    add_image(slides[11], FIGS / "fig2_tau_by_mechanism.png",
              left_in=0.3, top_in=1.2, width_in=7.0, height_in=3.8)
remove_shape_by_text(slides[11], "Kendall-? distance")
remove_shape_by_text(slides[11], "League fairness with and without honest")

# Slide 13 (index 12): mixed sweep
print("Slide 13: replacing image with fig3...")
if not replace_image(slides[12], FIGS / "fig3_mixed_sweep.png"):
    add_image(slides[12], FIGS / "fig3_mixed_sweep.png",
              left_in=0.3, top_in=1.2, width_in=7.0, height_in=3.8)
remove_shape_by_text(slides[12], "League tanking rates by number of rational")

# Slide 14 (index 13): LLM agents
print("Slide 14: replacing image and updating quote...")
if not replace_image(slides[13], llm_fig):
    add_image(slides[13], llm_fig,
              left_in=5.0, top_in=1.2, width_in=4.7, height_in=3.5)
remove_shape_by_text(slides[13], "League tanking rates by programmatic")
# Replace fabricated quote with real simulation output
update_text(
    slides[13],
    "Our current win probability is low enough",
    "At 14–18 we’re firmly out of playoff contention with 50 games left; "
    "minimizing effort maximizes losing to secure top-3 lottery odds."
)

# ── save ──────────────────────────────────────────────────────────────────────

prs.save(DEST)
print(f"\nSaved -> {DEST}")
